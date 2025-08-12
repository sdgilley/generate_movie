import os
import subprocess
import re
from pptx import Presentation
from moviepy import ImageClip, concatenate_videoclips
from PIL import Image
import tempfile
from dotenv import load_dotenv

# Import utility functions
try:
    from .filename_utils import get_powerpoint_file, generate_output_filename
except ImportError:
    from filename_utils import get_powerpoint_file, generate_output_filename

# Load environment variables
load_dotenv()

def export_slides_as_images_libreoffice(pptx_file, output_dir="exported_slides"):
    """Try to export slides using LibreOffice + ImageMagick with better visual fidelity"""
    os.makedirs(output_dir, exist_ok=True)
    
    try:
        # First, try to convert PowerPoint to PDF using LibreOffice with better settings
        print("Converting PowerPoint to PDF using LibreOffice (with enhanced settings)...")
        
        # Use more advanced LibreOffice export settings for better visual fidelity
        cmd_pdf = [
            "soffice", "--headless", "--invisible", "--convert-to", "pdf",
            "--outdir", ".",
            "--writer",  # Force through writer for better conversion
            pptx_file
        ]
        result_pdf = subprocess.run(cmd_pdf, capture_output=True, text=True, timeout=90)
        
        if result_pdf.returncode != 0:
            print(f"LibreOffice enhanced PDF conversion failed: {result_pdf.stderr}")
            # Try simpler conversion as fallback
            cmd_pdf_simple = [
                "soffice", "--headless", "--convert-to", "pdf",
                "--outdir", ".", pptx_file
            ]
            result_pdf = subprocess.run(cmd_pdf_simple, capture_output=True, text=True, timeout=60)
            
            if result_pdf.returncode != 0:
                print(f"LibreOffice simple PDF conversion also failed: {result_pdf.stderr}")
                return False
        
        # Check if PDF was created
        pdf_path = os.path.splitext(pptx_file)[0] + ".pdf"
        if not os.path.exists(pdf_path):
            print("PDF file was not created successfully")
            return False
        
        print(f"PDF created successfully: {pdf_path}")
        
        # Use ImageMagick with enhanced settings for better slide image quality
        print("Converting PDF pages to individual slide images using ImageMagick (high quality)...")
        cmd_convert = [
            "convert", 
            "-density", "300",      # Higher resolution for better quality
            "-quality", "95",       # Higher quality setting
            "-background", "white", # Ensure white background
            "-alpha", "remove",     # Remove transparency issues
            "-colorspace", "RGB",   # Ensure RGB colorspace
            pdf_path,
            os.path.join(output_dir, "slide_%d.png")
        ]
        
        result_convert = subprocess.run(cmd_convert, capture_output=True, text=True, timeout=180)
        
        # Clean up the temporary PDF
        if os.path.exists(pdf_path):
            os.remove(pdf_path)
            print(f"Cleaned up temporary PDF: {pdf_path}")
        
        if result_convert.returncode == 0:
            # Check how many slide images were created
            slide_files = [f for f in os.listdir(output_dir) if f.startswith('slide_') and f.endswith('.png')]
            print(f"Successfully exported {len(slide_files)} slides using LibreOffice + ImageMagick (enhanced)")
            return True
        else:
            print(f"ImageMagick conversion failed: {result_convert.stderr}")
            return False
            
    except Exception as e:
        print(f"LibreOffice + ImageMagick enhanced export failed: {e}")
        # Clean up PDF if it exists
        pdf_path = os.path.splitext(pptx_file)[0] + ".pdf"
        if os.path.exists(pdf_path):
            os.remove(pdf_path)
        return False

def export_slides_as_images_macos_keynote(pptx_file, output_dir="exported_slides"):
    """Try to export slides using macOS Keynote (if available)"""
    os.makedirs(output_dir, exist_ok=True)
    
    try:
        # Create AppleScript to open PowerPoint in Keynote and export as images
        applescript = f'''
        tell application "Keynote"
            activate
            open POSIX file "{os.path.abspath(pptx_file)}"
            delay 3
            tell front document
                set slide_count to count of slides
                repeat with i from 1 to slide_count
                    set slide_path to "{os.path.abspath(output_dir)}/slide_" & i & ".png"
                    export slide i as "slide images" to file slide_path
                end repeat
            end tell
            close front document
            quit
        end tell
        '''
        
        print("Trying to export slides using macOS Keynote...")
        result = subprocess.run(
            ["osascript", "-e", applescript],
            capture_output=True, text=True, timeout=120
        )
        
        if result.returncode == 0:
            slide_files = [f for f in os.listdir(output_dir) if f.startswith('slide_') and f.endswith('.png')]
            if len(slide_files) > 0:
                print(f"Successfully exported {len(slide_files)} slides using macOS Keynote")
                return True
            else:
                print("Keynote export completed but no slide files found")
                return False
        else:
            print(f"Keynote export failed: {result.stderr}")
            return False
            
    except Exception as e:
        print(f"macOS Keynote export failed: {e}")
        return False

def export_slides_as_images_powershell(pptx_file, output_dir="exported_slides"):
    """Try to export slides using PowerShell (Windows only)"""
    os.makedirs(output_dir, exist_ok=True)
    
    try:
        # PowerShell script to export PowerPoint slides
        ps_script = f'''
        $ppt = New-Object -ComObject PowerPoint.Application
        $ppt.Visible = [Microsoft.Office.Core.MsoTriState]::msoFalse
        $presentation = $ppt.Presentations.Open("{os.path.abspath(pptx_file)}")
        
        for ($i = 1; $i -le $presentation.Slides.Count; $i++) {{
            $slide = $presentation.Slides.Item($i)
            $slide.Export("{os.path.abspath(output_dir)}\\slide_$i.png", "PNG")
        }}
        
        $presentation.Close()
        $ppt.Quit()
        [System.Runtime.Interopservices.Marshal]::ReleaseComObject($ppt)
        '''
        
        # Save PowerShell script to temp file
        with tempfile.NamedTemporaryFile(mode='w', suffix='.ps1', delete=False) as f:
            f.write(ps_script)
            ps_file = f.name
        
        # Execute PowerShell script
        cmd = ["powershell.exe", "-ExecutionPolicy", "Bypass", "-File", ps_file]
        result = subprocess.run(cmd, capture_output=True, text=True, timeout=120)
        
        # Clean up temp file
        os.unlink(ps_file)
        
        if result.returncode == 0:
            print("Successfully exported slides using PowerShell")
            return True
        else:
            print(f"PowerShell export failed: {result.stderr}")
            return False
            
    except Exception as e:
        print(f"PowerShell export failed: {e}")
        return False
    """Try to export slides using PowerShell (Windows only)"""
    os.makedirs(output_dir, exist_ok=True)
    
    try:
        # PowerShell script to export PowerPoint slides
        ps_script = f'''
        $ppt = New-Object -ComObject PowerPoint.Application
        $ppt.Visible = [Microsoft.Office.Core.MsoTriState]::msoFalse
        $presentation = $ppt.Presentations.Open("{os.path.abspath(pptx_file)}")
        
        for ($i = 1; $i -le $presentation.Slides.Count; $i++) {{
            $slide = $presentation.Slides.Item($i)
            $slide.Export("{os.path.abspath(output_dir)}\\slide_$i.png", "PNG")
        }}
        
        $presentation.Close()
        $ppt.Quit()
        [System.Runtime.Interopservices.Marshal]::ReleaseComObject($ppt)
        '''
        
        # Save PowerShell script to temp file
        with tempfile.NamedTemporaryFile(mode='w', suffix='.ps1', delete=False) as f:
            f.write(ps_script)
            ps_file = f.name
        
        # Execute PowerShell script
        cmd = ["powershell.exe", "-ExecutionPolicy", "Bypass", "-File", ps_file]
        result = subprocess.run(cmd, capture_output=True, text=True, timeout=120)
        
        # Clean up temp file
        os.unlink(ps_file)
        
        if result.returncode == 0:
            print("Successfully exported slides using PowerShell")
            return True
        else:
            print(f"PowerShell export failed: {result.stderr}")
            return False
            
    except Exception as e:
        print(f"PowerShell export failed: {e}")
        return False

def export_slides_python_fallback(pptx_file, output_dir="exported_slides"):
    """Export slides using pure Python approach (macOS compatible)"""
    print("Trying Python-based slide export (macOS compatible)...")
    
    try:
        from PIL import Image, ImageDraw, ImageFont
        
        def extract_slide_content(slide):
            """Extract text and basic layout from a slide"""
            content = {
                'title': '',
                'bullet_points': [],
                'other_text': []
            }
            
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    text = shape.text.strip()
                    
                    # Try to identify title (usually first or short text)
                    try:
                        if hasattr(shape, "placeholder_format") and shape.placeholder_format:
                            if shape.placeholder_format.type == 1:  # Title placeholder
                                content['title'] = text
                            elif shape.placeholder_format.type == 2:  # Content placeholder
                                # Split by lines and treat as bullet points
                                lines = [line.strip() for line in text.split('\n') if line.strip()]
                                content['bullet_points'].extend(lines)
                            else:
                                content['other_text'].append(text)
                        else:
                            # No placeholder info, make educated guess
                            if not content['title'] and len(text) < 100:
                                content['title'] = text
                            else:
                                content['other_text'].append(text)
                    except Exception as e:
                        # If placeholder access fails, treat as other text
                        print(f"Warning: Could not access placeholder format for shape: {e}")
                        if not content['title'] and len(text) < 100:
                            content['title'] = text
                        else:
                            content['other_text'].append(text)
            
            return content

        def create_slide_image(slide_content, slide_num, width=1280, height=720):
            """Create a simple slide image from extracted content"""
            # Create a white background
            img = Image.new('RGB', (width, height), color='white')
            draw = ImageDraw.Draw(img)
            
            # Try to load system fonts, fallback to default
            try:
                # macOS system fonts
                title_font = ImageFont.truetype('/System/Library/Fonts/Helvetica.ttc', 48)
                content_font = ImageFont.truetype('/System/Library/Fonts/Helvetica.ttc', 32)
                small_font = ImageFont.truetype('/System/Library/Fonts/Helvetica.ttc', 24)
            except:
                try:
                    # Fallback fonts
                    title_font = ImageFont.truetype('Arial.ttf', 48)
                    content_font = ImageFont.truetype('Arial.ttf', 32)
                    small_font = ImageFont.truetype('Arial.ttf', 24)
                except:
                    # Use default font as last resort
                    title_font = ImageFont.load_default()
                    content_font = ImageFont.load_default()
                    small_font = ImageFont.load_default()
            
            y_position = 80
            margin = 60
            
            # Draw title
            if slide_content['title']:
                # Center the title
                title_bbox = draw.textbbox((0, 0), slide_content['title'], font=title_font)
                title_width = title_bbox[2] - title_bbox[0]
                title_x = (width - title_width) // 2
                draw.text((title_x, y_position), slide_content['title'], fill='black', font=title_font)
                y_position += 100
            
            # Draw bullet points
            for bullet in slide_content['bullet_points']:
                if y_position > height - 100:  # Leave space at bottom
                    break
                # Wrap long text
                words = bullet.split()
                lines = []
                current_line = []
                
                for word in words:
                    test_line = ' '.join(current_line + [word])
                    bbox = draw.textbbox((0, 0), test_line, font=content_font)
                    if bbox[2] - bbox[0] < width - 2 * margin:
                        current_line.append(word)
                    else:
                        if current_line:
                            lines.append(' '.join(current_line))
                        current_line = [word]
                
                if current_line:
                    lines.append(' '.join(current_line))
                
                for line in lines:
                    draw.text((margin, y_position), f"• {line}", fill='black', font=content_font)
                    y_position += 45
            
            # Draw other text
            for text in slide_content['other_text']:
                if y_position > height - 100:
                    break
                # Simple text wrapping
                words = text.split()
                lines = []
                current_line = []
                
                for word in words:
                    test_line = ' '.join(current_line + [word])
                    bbox = draw.textbbox((0, 0), test_line, font=small_font)
                    if bbox[2] - bbox[0] < width - 2 * margin:
                        current_line.append(word)
                    else:
                        if current_line:
                            lines.append(' '.join(current_line))
                        current_line = [word]
                
                if current_line:
                    lines.append(' '.join(current_line))
                
                for line in lines:
                    draw.text((margin, y_position), line, fill='gray', font=small_font)
                    y_position += 35
            
            # Add slide number
            slide_text = f"Slide {slide_num}"
            draw.text((width - 150, height - 40), slide_text, fill='gray', font=small_font)
            
            return img

        # Load presentation
        presentation = Presentation(pptx_file)
        print(f"Loaded presentation with {len(presentation.slides)} slides for Python export")
        
        # Create output directory
        os.makedirs(output_dir, exist_ok=True)
        
        # Process each slide
        for i, slide in enumerate(presentation.slides):
            slide_num = i + 1
            print(f"Creating slide image {slide_num}...")
            
            try:
                # Extract content
                content = extract_slide_content(slide)
                
                # Create image
                img = create_slide_image(content, slide_num)
                
                # Save image
                output_path = os.path.join(output_dir, f"slide_{slide_num}.png")
                img.save(output_path)
            except Exception as e:
                print(f"Warning: Failed to process slide {slide_num}: {e}")
                # Create a simple error slide
                try:
                    error_img = Image.new('RGB', (1280, 720), color='white')
                    error_draw = ImageDraw.Draw(error_img)
                    error_draw.text((50, 300), f"Slide {slide_num} - Processing Error", fill='red')
                    output_path = os.path.join(output_dir, f"slide_{slide_num}.png")
                    error_img.save(output_path)
                except:
                    print(f"Could not create error slide for slide {slide_num}")
                    continue
        
        print(f"Successfully exported {len(presentation.slides)} slides to {output_dir}/")
        return True
        
    except Exception as e:
        print(f"Python-based export failed: {e}")
        return False

def export_slides_manual_instruction(pptx_file):
    """Provide manual instructions for exporting slides"""
    print(f"""
MANUAL EXPORT INSTRUCTIONS:
1. Open {pptx_file} in PowerPoint
2. Go to File > Export > Change File Type
3. Select PNG Portable Network Graphics Format
4. Click Save As
5. Choose a folder name like 'manual_slides'
6. PowerPoint will ask if you want to export "All Slides" - click "All Slides"
7. This will create individual PNG files for each slide
8. Come back and run this script again, it will look for the exported images
""")

def find_exported_slides():
    """Look for exported slide images in common locations"""
    possible_dirs = ["uploaded_slides", "exported_slides", "manual_slides", "slide_exports", "slides"]
    
    for dir_name in possible_dirs:
        if os.path.exists(dir_name):
            # Accept both slide_1.png and Slide1.png patterns
            png_files = [f for f in os.listdir(dir_name) if f.lower().endswith('.png') and (f.lower().startswith('slide_') or f.lower().startswith('slide'))]
            if png_files:
                # Sort numerically instead of alphabetically
                import re
                def natural_sort_key(filename):
                    # Extract numbers from filename for proper sorting
                    numbers = re.findall(r'\d+', filename)
                    if numbers:
                        return int(numbers[0])  # Sort by first number found
                    return 0
                png_files.sort(key=natural_sort_key)
                return dir_name, png_files
    
    return None, []

def extract_narration_from_slides(presentation):
    """Extract narration text from PowerPoint slide notes"""
    narration_dict = {}
    
    for i, slide in enumerate(presentation.slides):
        try:
            # Get slide title if available
            title = ""
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    if hasattr(shape, "placeholder_format") and shape.placeholder_format.type == 1:  # Title placeholder
                        title = shape.text.strip()
                        break
                    elif not title and len(shape.text.strip()) < 100:  # Fallback: short text might be title
                        title = shape.text.strip()
            
            # Get slide notes
            if slide.has_notes_slide:
                notes_slide = slide.notes_slide
                if notes_slide.notes_text_frame:
                    narration_text = notes_slide.notes_text_frame.text.strip()
                    if narration_text:
                        # Store by title and by slide number as fallback
                        if title:
                            print(f"Slide {i+1} ('{title}'): {len(narration_text)} characters of narration")
                            narration_dict[title] = narration_text
                        
                        # Also store by slide number for backup lookup
                        slide_key = f"slide_{i+1}"
                        narration_dict[slide_key] = narration_text
                        
                        if not title:
                            print(f"Slide {i+1} (no title): {len(narration_text)} characters of narration")
        
        except Exception as e:
            print(f"Warning: Could not extract narration from slide {i+1}: {e}")
            continue
    
    return narration_dict

# Load the PowerPoint presentation to get slide count and titles
def main():
    """Main function to export slides from PowerPoint"""
    pptx_file = os.environ.get('POWERPOINT_FILE', 'content_maintenance_process.pptx')
    print(f"Loading PowerPoint presentation: {pptx_file}")
    presentation = Presentation(pptx_file)
    print(f"Loaded presentation with {len(presentation.slides)} slides")

    # Extract narration from PowerPoint slide notes
    print(f"\nExtracting narration from PowerPoint slide notes...")
    narration_notes = extract_narration_from_slides(presentation)

    if not narration_notes:
        print("Warning: No narration found in slide notes. Videos will be silent.")
    else:
        print(f"Ready to process with {len(narration_notes)} narration entries")

    # Try to export slides automatically
    print("\nAttempting to export slides as images...")

    # If uploaded_slides exists and contains PNGs, use those and skip export/cleanup
    uploaded_dir = "uploaded_slides"
    uploaded_pngs = [f for f in os.listdir(uploaded_dir) if f.lower().endswith('.png')] if os.path.exists(uploaded_dir) else []
    if uploaded_pngs:
        print(f"Found {len(uploaded_pngs)} PNG slides in {uploaded_dir}/. Using these and skipping export.")
        slides_dir, slide_files = uploaded_dir, uploaded_pngs
    else:
        # Always clean up exported_slides before generating new ones
        export_dir = "exported_slides"
        if os.path.exists(export_dir):
            old_files = [f for f in os.listdir(export_dir) if f.lower().endswith('.png')]
            if old_files:
                print(f"Cleaning up {len(old_files)} old slide images from {export_dir}/.")
                for f in old_files:
                    os.remove(os.path.join(export_dir, f))

        success = False

        # Try macOS Keynote method first (best visual fidelity on macOS)
        print("Trying macOS Keynote export...")
        success = export_slides_as_images_macos_keynote(pptx_file)

        # Try PowerShell method (best for Windows)
        if not success:
            print("Trying PowerShell export...")
            success = export_slides_as_images_powershell(pptx_file)

        # If PowerShell failed, try LibreOffice
        if not success:
            print("Trying LibreOffice export...")
            success = export_slides_as_images_libreoffice(pptx_file)

        # If both failed, try Python-based approach (macOS compatible)
        if not success:
            print("Trying Python-based export (macOS compatible)...")
            success = export_slides_python_fallback(pptx_file)

        # Look for existing exported slides
        slides_dir, slide_files = find_exported_slides()

    # Look for existing exported slides
    slides_dir, slide_files = find_exported_slides()

    if not slides_dir:
        print("Could not automatically export slides.")
        export_slides_manual_instruction(pptx_file)
        print("\nAfter manually exporting, run this script again.")
        return False

    print(f"Found {len(slide_files)} exported slide images in {slides_dir}/")
    for i, filename in enumerate(slide_files[:5]):  # Show first 5
        print(f"  {i+1}: {filename}")
    if len(slide_files) > 5:
        print(f"  ... and {len(slide_files)-5} more")

    # Create video from exported slides
    video_clips = []
    print("\nCreating video from slide images...")

    for i, slide_file in enumerate(slide_files):
        slide_path = os.path.join(slides_dir, slide_file)
        print(f"\n--- Processing slide {i+1}/{len(slide_files)}: {slide_file} ---")
        
        try:
            # Load the slide image
            slide_img = Image.open(slide_path)
            print(f"Loaded slide image: {slide_img.size}")
            
            # Resize if needed to standard video dimensions
            # Keep aspect ratio but fit within 1280x720
            slide_img.thumbnail((1280, 720), Image.Resampling.LANCZOS)
            
            # Create a 1280x720 canvas and center the slide
            canvas = Image.new("RGB", (1280, 720), color="white")
            x_offset = (1280 - slide_img.width) // 2
            y_offset = (720 - slide_img.height) // 2
            canvas.paste(slide_img, (x_offset, y_offset))
            
            # Save the processed slide
            processed_path = f"slide_images/processed_slide_{i+1}.png"
            os.makedirs("slide_images", exist_ok=True)
            canvas.save(processed_path)
            
            # Create video clip
            duration = 4  # 4 seconds per slide
            image_clip = ImageClip(processed_path).with_duration(duration)
            video_clips.append(image_clip)
            print(f"Created video clip {i+1} with duration {duration}s")
            
        except Exception as e:
            print(f"Error processing slide {i+1}: {e}")

    # Combine clips into final video
    if video_clips:
        print(f"\nCombining {len(video_clips)} video clips into final video...")
        final_video = concatenate_videoclips(video_clips, method="compose")
        print("Video clips concatenated successfully")

        # Generate output filename based on PowerPoint file (for intermediate slides-only video)
        pptx_file = get_powerpoint_file()
        output_filename = generate_output_filename(pptx_file, "_SLIDES")
        print(f"Writing final video file: {output_filename}")
        final_video.write_videofile(output_filename, fps=24)
        print("Final video file written successfully!")

        # Cleanup
        print("Starting cleanup...")
        for clip in video_clips:
            clip.close()
        final_video.close()
        
        print(f"\nVideo created successfully: {output_filename}")
        print(f"Duration: {len(video_clips) * 4} seconds ({len(video_clips)} slides x 4 seconds each)")
        
        # Return the filename so it can be cleaned up later if needed
        return output_filename
    else:
        print("No video clips were created!")
        return None

if __name__ == "__main__":
    main()
