import os
import re
import subprocess
import numpy as np
from pptx import Presentation
from moviepy import ImageClip, AudioFileClip, concatenate_videoclips
from PIL import Image
import tempfile
from dotenv import load_dotenv

# Import generate_audio_file with fallback for relative import
try:
    from .generate_audio import generate_audio_file
    from .filename_utils import get_powerpoint_file, get_output_video_name
except ImportError:
    from generate_audio import generate_audio_file
    from filename_utils import get_powerpoint_file, get_output_video_name

# Load environment variables
load_dotenv()

# Configuration from environment variables
pptx_file = get_powerpoint_file()
output_video_name = get_output_video_name()
pause_duration = float(os.environ.get('PAUSE_DURATION', '1.5'))
include_end_slide = os.environ.get('INCLUDE_END_SLIDE', 'true').lower() in ('true', 'yes', '1')

def extract_narration_from_slides(presentation):
    """Extract narration text from PowerPoint slide notes"""
    narration_dict = {}
    
    for i, slide in enumerate(presentation.slides):
        try:
            # Get slide title if available
            title = ""
            if hasattr(slide, 'shapes') and slide.shapes.title:
                title = slide.shapes.title.text.strip()
            
            # Get narration from slide notes
            narration = ""
            try:
                notes_slide = slide.notes_slide
                if notes_slide and hasattr(notes_slide, 'notes_text_frame'):
                    notes_text_frame = notes_slide.notes_text_frame
                    if notes_text_frame and hasattr(notes_text_frame, 'text'):
                        narration = notes_text_frame.text.strip()
            except Exception as e:
                print(f"Warning: Could not read notes for slide {i+1}: {e}")
            
            # Store narration by title (for compatibility) and by slide number
            if title and narration:
                narration_dict[title] = narration
                narration_dict[f"slide_{i+1}"] = narration
                print(f"Slide {i+1} ('{title}'): {len(narration)} characters of narration")
            elif narration:
                narration_dict[f"slide_{i+1}"] = narration
                print(f"Slide {i+1} (no title): {len(narration)} characters of narration")
            else:
                print(f"Slide {i+1}: No narration found")
                
        except Exception as e:
            print(f"Error processing slide {i+1}: {e}")
    
    return narration_dict

def find_exported_slides():
    """Look for exported slide images in common locations"""
    possible_dirs = ["uploaded_slides", "exported_slides", "manual_slides", "slide_exports", "slides"]
    
    for dir_name in possible_dirs:
        if os.path.exists(dir_name):
            png_files = [f for f in os.listdir(dir_name) if f.lower().endswith('.png')]
            if png_files:
                # Sort numerically instead of alphabetically
                def natural_sort_key(filename):
                    # Extract numbers from filename for proper sorting
                    numbers = re.findall(r'\d+', filename)
                    if numbers:
                        return int(numbers[0])  # Sort by first number found
                    return 0
                
                png_files.sort(key=natural_sort_key)
                return dir_name, png_files
    
    return None, []

def main():
    """Main function to convert PowerPoint to video with Azure Speech narration"""
    try:
        print(f"Configuration loaded from .env:")
        print(f"  PowerPoint file: {pptx_file}")
        print(f"  Output video: {output_video_name}")
        print(f"  Pause duration: {pause_duration}s")
        print(f"  Include end slide: {include_end_slide}")

        # Load the PowerPoint presentation
        print(f"\nLoading PowerPoint presentation: {pptx_file}")
        if not os.path.exists(pptx_file):
            print(f"ERROR: PowerPoint file not found: {pptx_file}")
            print("Please check the POWERPOINT_FILE setting in your .env file")
            return False

        presentation = Presentation(pptx_file)
        print(f"Loaded presentation with {len(presentation.slides)} slides")

        # Extract narration from PowerPoint slide notes
        print(f"\nExtracting narration from PowerPoint slide notes...")
        narration_notes = extract_narration_from_slides(presentation)

        if not narration_notes:
            print("Warning: No narration found in slide notes. Videos will be silent.")
        else:
            print(f"Ready to process with {len(narration_notes)} narration entries")

        # Look for existing exported slides
        slides_dir, slide_files = find_exported_slides()

        if not slides_dir:
            print("No exported slides found. Please run generate_from_slides.py first to export slides.")
            return False

        print(f"Found {len(slide_files)} exported slide images in {slides_dir}/")

        # Create directories for audio and final processing
        os.makedirs("audio_clips", exist_ok=True)
        os.makedirs("slide_images", exist_ok=True)

        # Process slides and generate audio
        video_clips = []
        print("\nCreating video with Azure Speech Services audio...")

        for i, slide_file in enumerate(slide_files):
            slide_path = os.path.join(slides_dir, slide_file)
            print(f"\n--- Processing slide {i+1}/{len(slide_files)}: {slide_file} ---")
            
            try:
                # Load the slide image
                slide_img = Image.open(slide_path)
                print(f"Loaded slide image: {slide_img.size}")
                
                # Resize if needed to standard video dimensions
                slide_img.thumbnail((1280, 720), Image.Resampling.LANCZOS)
                
                # Create a 1280x720 canvas and center the slide
                canvas = Image.new("RGB", (1280, 720), color="white")
                x_offset = (1280 - slide_img.width) // 2
                y_offset = (720 - slide_img.height) // 2
                canvas.paste(slide_img, (x_offset, y_offset))
                
                # Save the processed slide
                processed_path = f"slide_images/processed_slide_{i+1}.png"
                canvas.save(processed_path)
                
                # Get slide title and narration
                slide = presentation.slides[i] if i < len(presentation.slides) else None
                title = ""
                narration = ""
                
                if slide and hasattr(slide, 'shapes') and slide.shapes.title:
                    title = slide.shapes.title.text.strip()
                
                # Try multiple ways to find narration:
                # 1. By slide title (for compatibility with external files)
                # 2. By slide number (slide_1, slide_2, etc.)
                # 3. Direct lookup in case title is not available
                if title:
                    narration = narration_notes.get(title, "")
                    if narration:
                        print(f"Found narration for '{title}' by title: {len(narration)} characters")
                
                if not narration:
                    slide_key = f"slide_{i+1}"
                    narration = narration_notes.get(slide_key, "")
                    if narration:
                        print(f"Found narration for slide {i+1} by number: {len(narration)} characters")
                
                if not narration:
                    print(f"No narration found for slide {i+1} (title: '{title}')")
                
                if narration:
                    # Generate audio using Azure Speech Services
                    audio_path = f"audio_clips/audio_{i+1}.wav"
                    audio_success = generate_audio_file(narration, audio_path)
                    
                    if audio_success and os.path.exists(audio_path):
                        # Create slide with pause + narration
                        audio_clip = AudioFileClip(audio_path)
                        audio_duration = audio_clip.duration
                        print(f"Audio duration: {audio_duration:.1f} seconds")
                        
                        # Create silent pause at the beginning (for reading the slide)
                        silent_clip = ImageClip(processed_path).with_duration(pause_duration)
                        print(f"Added {pause_duration}s reading pause before narration")
                        
                        # Create narrated portion
                        narrated_clip = ImageClip(processed_path).with_duration(audio_duration).with_audio(audio_clip)
                        
                        # Combine pause + narration for this slide
                        from moviepy import concatenate_videoclips
                        slide_video = concatenate_videoclips([silent_clip, narrated_clip])
                        
                        total_duration = pause_duration + audio_duration
                        video_clips.append(slide_video)
                        print(f"Created slide {i+1}: {pause_duration}s pause + {audio_duration:.1f}s narration = {total_duration:.1f}s total")
                    else:
                        # Fallback: create clip without audio
                        print("Audio generation failed, creating silent clip")
                        duration = 5  # 5 seconds default
                        image_clip = ImageClip(processed_path).with_duration(duration)
                        video_clips.append(image_clip)
                else:
                    # No narration found, create short silent clip
                    print(f"No narration found for slide {i+1}, creating short silent clip")
                    duration = 3  # 3 seconds for slides without narration
                    image_clip = ImageClip(processed_path).with_duration(duration)
                    video_clips.append(image_clip)
                
            except Exception as e:
                print(f"Error processing slide {i+1}: {e}")

        # Combine clips into final video
        if video_clips:
            print(f"\nCombining {len(video_clips)} slide clips into final video...")
            print("Each slide now includes its own reading pause before narration")
            
            # Add end slide (if enabled)
            if include_end_slide:
                print("Adding end slide...")
                try:
                    # Use the existing end slide image file
                    end_slide_path = "media/end-slide.png"
                    
                    # Verify the file exists
                    if not os.path.exists(end_slide_path):
                        raise FileNotFoundError(f"End slide image not found: {end_slide_path}")
                    
                    print(f"Using end slide image: {end_slide_path}")
                    
                    # Generate narration for the end slide
                    end_slide_narration = "This video was generated from a PowerPoint presentation using Azure AI. For more information and source code, visit github.com/sdgilley/generate_movie"
                    end_slide_audio_path = "audio_clips/audio_end_slide.wav"
                    
                    print("Generating audio for end slide...")
                    audio_success = generate_audio_file(end_slide_narration, end_slide_audio_path)
                    
                    if audio_success and os.path.exists(end_slide_audio_path):
                        # Create end slide with audio
                        audio_clip = AudioFileClip(end_slide_audio_path)
                        audio_duration = audio_clip.duration
                        print(f"End slide audio duration: {audio_duration:.1f} seconds")
                        
                        # Create silent pause at the beginning (for reading the slide)
                        silent_end_clip = ImageClip(end_slide_path).with_duration(pause_duration)
                        print(f"Added {pause_duration}s reading pause before end slide narration")
                        
                        # Create narrated portion
                        narrated_end_clip = ImageClip(end_slide_path).with_duration(audio_duration).with_audio(audio_clip)
                        
                        # Combine pause + narration for end slide
                        end_slide_video = concatenate_videoclips([silent_end_clip, narrated_end_clip])
                        
                        total_end_duration = pause_duration + audio_duration
                        video_clips.append(end_slide_video)
                        print(f"End slide added to video ({pause_duration}s pause + {audio_duration:.1f}s narration = {total_end_duration:.1f}s total)")
                    else:
                        # Fallback: create silent end slide (5 seconds duration)
                        print("Warning: Could not generate end slide audio, using silent version")
                        end_slide_clip = ImageClip(end_slide_path).with_duration(5.0)
                        video_clips.append(end_slide_clip)
                        print("End slide added to video (5 seconds, silent)")
                    
                except Exception as e:
                    print(f"Warning: Could not create end slide: {e}")
            else:
                print("End slide skipped (INCLUDE_END_SLIDE=false)")
            
            # Combine all slide clips (no additional pauses needed)
            final_video = concatenate_videoclips(video_clips, method="compose")
            print("Video clips concatenated successfully")

            print(f"Writing final video file: {output_video_name}")
            final_video.write_videofile(output_video_name, fps=24)
            print("Final video file written successfully!")

            # Calculate total duration
            total_duration = sum([clip.duration for clip in video_clips])
            
            print(f"\nVideo created successfully: {output_video_name}")
            print(f"Total duration: {total_duration:.1f} seconds ({total_duration/60:.1f} minutes)")
            print(f"Number of slides: {len(slide_files)}")

            # Cleanup
            print("Starting cleanup...")
            # Close video clips to free memory
            for clip in video_clips:
                try:
                    if hasattr(clip, 'close'):
                        clip.close()
                except:
                    pass  # Ignore any cleanup errors
            
            try:
                final_video.close()
            except:
                pass
            
            # Clean up intermediate files
            print("Cleaning up intermediate files...")
            
            # Remove processed slide images (we keep the original exported slides)
            if os.path.exists("slide_images"):
                for file in os.listdir("slide_images"):
                    file_path = os.path.join("slide_images", file)
                    os.remove(file_path)
                    print(f"Removed: {file_path}")
                os.rmdir("slide_images")
                print("Removed directory: slide_images/")
            
            # Keep audio files for now, user might want them
            # If you want to remove audio files too, uncomment below:
            # if os.path.exists("audio_clips"):
            #     for file in os.listdir("audio_clips"):
            #         file_path = os.path.join("audio_clips", file)
            #         os.remove(file_path)
            #         print(f"Removed: {file_path}")
            #     os.rmdir("audio_clips")
            #     print("Removed directory: audio_clips/")
            
            # Remove test audio if it exists
            if os.path.exists("test_audio"):
                for file in os.listdir("test_audio"):
                    file_path = os.path.join("test_audio", file)
                    os.remove(file_path)
                    print(f"Removed: {file_path}")
                os.rmdir("test_audio")
                print("Removed directory: test_audio/")
            
            print("Cleanup completed!")
            
        else:
            print("No video clips were created!")
            return False

        return True
        
    except Exception as e:
        print(f"Error in video generation: {e}")
        return False

if __name__ == "__main__":
    main()