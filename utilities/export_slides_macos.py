#!/usr/bin/env python3
"""
macOS-compatible PowerPoint slide export utility
Uses python-pptx to extract slide content and create images
"""

import os
import sys
from pptx import Presentation
from PIL import Image, ImageDraw, ImageFont
import io
from dotenv import load_dotenv

load_dotenv()

def extract_slide_content(slide):
    """Extract text and basic layout from a slide"""
    content = {
        'title': '',
        'bullet_points': [],
        'other_text': []
    }
    
    for shape in slide.shapes:
        if hasattr(shape, "text") and shape.text.strip():
            text = shape.text.strip()
            
            # Try to identify title (usually first or short text)
            if hasattr(shape, "placeholder_format"):
                if shape.placeholder_format.type == 1:  # Title placeholder
                    content['title'] = text
                elif shape.placeholder_format.type == 2:  # Content placeholder
                    # Split by lines and treat as bullet points
                    lines = [line.strip() for line in text.split('\n') if line.strip()]
                    content['bullet_points'].extend(lines)
                else:
                    content['other_text'].append(text)
            else:
                # No placeholder info, make educated guess
                if not content['title'] and len(text) < 100:
                    content['title'] = text
                else:
                    content['other_text'].append(text)
    
    return content

def create_slide_image(slide_content, slide_num, width=1280, height=720):
    """Create a simple slide image from extracted content"""
    # Create a white background
    img = Image.new('RGB', (width, height), color='white')
    draw = ImageDraw.Draw(img)
    
    # Try to load system fonts, fallback to default
    try:
        # macOS system fonts
        title_font = ImageFont.truetype('/System/Library/Fonts/Helvetica.ttc', 48)
        content_font = ImageFont.truetype('/System/Library/Fonts/Helvetica.ttc', 32)
        small_font = ImageFont.truetype('/System/Library/Fonts/Helvetica.ttc', 24)
    except:
        try:
            # Fallback fonts
            title_font = ImageFont.truetype('Arial.ttf', 48)
            content_font = ImageFont.truetype('Arial.ttf', 32)
            small_font = ImageFont.truetype('Arial.ttf', 24)
        except:
            # Use default font as last resort
            title_font = ImageFont.load_default()
            content_font = ImageFont.load_default()
            small_font = ImageFont.load_default()
    
    y_position = 80
    margin = 60
    
    # Draw title
    if slide_content['title']:
        # Center the title
        title_bbox = draw.textbbox((0, 0), slide_content['title'], font=title_font)
        title_width = title_bbox[2] - title_bbox[0]
        title_x = (width - title_width) // 2
        draw.text((title_x, y_position), slide_content['title'], fill='black', font=title_font)
        y_position += 100
    
    # Draw bullet points
    for bullet in slide_content['bullet_points']:
        if y_position > height - 100:  # Leave space at bottom
            break
        # Wrap long text
        words = bullet.split()
        lines = []
        current_line = []
        
        for word in words:
            test_line = ' '.join(current_line + [word])
            bbox = draw.textbbox((0, 0), test_line, font=content_font)
            if bbox[2] - bbox[0] < width - 2 * margin:
                current_line.append(word)
            else:
                if current_line:
                    lines.append(' '.join(current_line))
                current_line = [word]
        
        if current_line:
            lines.append(' '.join(current_line))
        
        for line in lines:
            draw.text((margin, y_position), f"• {line}", fill='black', font=content_font)
            y_position += 45
    
    # Draw other text
    for text in slide_content['other_text']:
        if y_position > height - 100:
            break
        # Simple text wrapping
        words = text.split()
        lines = []
        current_line = []
        
        for word in words:
            test_line = ' '.join(current_line + [word])
            bbox = draw.textbbox((0, 0), test_line, font=small_font)
            if bbox[2] - bbox[0] < width - 2 * margin:
                current_line.append(word)
            else:
                if current_line:
                    lines.append(' '.join(current_line))
                current_line = [word]
        
        if current_line:
            lines.append(' '.join(current_line))
        
        for line in lines:
            draw.text((margin, y_position), line, fill='gray', font=small_font)
            y_position += 35
    
    # Add slide number
    slide_text = f"Slide {slide_num}"
    draw.text((width - 150, height - 40), slide_text, fill='gray', font=small_font)
    
    return img

def export_slides_python(pptx_file, output_dir="exported_slides"):
    """Export slides using pure Python approach"""
    print(f"Attempting Python-based slide export for: {pptx_file}")
    
    try:
        # Load presentation
        presentation = Presentation(pptx_file)
        print(f"Loaded presentation with {len(presentation.slides)} slides")
        
        # Create output directory
        os.makedirs(output_dir, exist_ok=True)
        
        # Process each slide
        for i, slide in enumerate(presentation.slides):
            slide_num = i + 1
            print(f"Processing slide {slide_num}...")
            
            # Extract content
            content = extract_slide_content(slide)
            
            # Create image
            img = create_slide_image(content, slide_num)
            
            # Save image
            output_path = os.path.join(output_dir, f"slide_{slide_num}.png")
            img.save(output_path)
            print(f"Saved: {output_path}")
        
        print(f"Successfully exported {len(presentation.slides)} slides to {output_dir}/")
        return True
        
    except Exception as e:
        print(f"Error in Python-based export: {e}")
        return False

def main():
    """Test the export function"""
    pptx_file = os.environ.get('POWERPOINT_FILE', 'test-ppt.pptx')
    success = export_slides_python(pptx_file)
    return success

if __name__ == "__main__":
    main()
