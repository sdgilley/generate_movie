import os
import subprocess
import re
from pptx import Presentation
from moviepy import ImageClip, concatenate_videoclips
from PIL import Image
import tempfile

def export_slides_as_images_libreoffice(pptx_file, output_dir="exported_slides"):
    """Try to export slides using LibreOffice (if available)"""
    os.makedirs(output_dir, exist_ok=True)
    
    try:
        # Try LibreOffice command
        cmd = [
            "soffice", "--headless", "--convert-to", "png", 
            "--outdir", output_dir, pptx_file
        ]
        result = subprocess.run(cmd, capture_output=True, text=True, timeout=60)
        
        if result.returncode == 0:
            print("Successfully exported slides using LibreOffice")
            return True
        else:
            print(f"LibreOffice export failed: {result.stderr}")
            return False
    except Exception as e:
        print(f"LibreOffice not available or failed: {e}")
        return False

def export_slides_as_images_powershell(pptx_file, output_dir="exported_slides"):
    """Try to export slides using PowerShell (Windows only)"""
    os.makedirs(output_dir, exist_ok=True)
    
    try:
        # PowerShell script to export PowerPoint slides
        ps_script = f'''
        $ppt = New-Object -ComObject PowerPoint.Application
        $ppt.Visible = [Microsoft.Office.Core.MsoTriState]::msoFalse
        $presentation = $ppt.Presentations.Open("{os.path.abspath(pptx_file)}")
        
        for ($i = 1; $i -le $presentation.Slides.Count; $i++) {{
            $slide = $presentation.Slides.Item($i)
            $slide.Export("{os.path.abspath(output_dir)}\\slide_$i.png", "PNG")
        }}
        
        $presentation.Close()
        $ppt.Quit()
        [System.Runtime.Interopservices.Marshal]::ReleaseComObject($ppt)
        '''
        
        # Save PowerShell script to temp file
        with tempfile.NamedTemporaryFile(mode='w', suffix='.ps1', delete=False) as f:
            f.write(ps_script)
            ps_file = f.name
        
        # Execute PowerShell script
        cmd = ["powershell.exe", "-ExecutionPolicy", "Bypass", "-File", ps_file]
        result = subprocess.run(cmd, capture_output=True, text=True, timeout=120)
        
        # Clean up temp file
        os.unlink(ps_file)
        
        if result.returncode == 0:
            print("Successfully exported slides using PowerShell")
            return True
        else:
            print(f"PowerShell export failed: {result.stderr}")
            return False
            
    except Exception as e:
        print(f"PowerShell export failed: {e}")
        return False

def export_slides_manual_instruction(pptx_file):
    """Provide manual instructions for exporting slides"""
    print(f"""
MANUAL EXPORT INSTRUCTIONS:
1. Open {pptx_file} in PowerPoint
2. Go to File > Export > Change File Type
3. Select PNG Portable Network Graphics Format
4. Click Save As
5. Choose a folder name like 'manual_slides'
6. PowerPoint will ask if you want to export "All Slides" - click "All Slides"
7. This will create individual PNG files for each slide
8. Come back and run this script again, it will look for the exported images
""")

def find_exported_slides():
    """Look for exported slide images in common locations"""
    possible_dirs = ["exported_slides", "manual_slides", "slide_exports", "slides"]
    
    for dir_name in possible_dirs:
        if os.path.exists(dir_name):
            png_files = [f for f in os.listdir(dir_name) if f.lower().endswith('.png')]
            if png_files:
                # Sort numerically instead of alphabetically
                import re
                def natural_sort_key(filename):
                    # Extract numbers from filename for proper sorting
                    numbers = re.findall(r'\d+', filename)
                    if numbers:
                        return int(numbers[0])  # Sort by first number found
                    return 0
                
                png_files.sort(key=natural_sort_key)
                return dir_name, png_files
    
    return None, []

def extract_narration_from_slides(presentation):
    """Extract narration text from PowerPoint slide notes"""
    narration_dict = {}
    
    for i, slide in enumerate(presentation.slides):
        try:
            # Get slide title if available
            title = ""
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    if hasattr(shape, "placeholder_format") and shape.placeholder_format.type == 1:  # Title placeholder
                        title = shape.text.strip()
                        break
                    elif not title and len(shape.text.strip()) < 100:  # Fallback: short text might be title
                        title = shape.text.strip()
            
            # Get slide notes
            if slide.has_notes_slide:
                notes_slide = slide.notes_slide
                if notes_slide.notes_text_frame:
                    narration_text = notes_slide.notes_text_frame.text.strip()
                    if narration_text:
                        # Store by title and by slide number as fallback
                        if title:
                            print(f"Slide {i+1} ('{title}'): {len(narration_text)} characters of narration")
                            narration_dict[title] = narration_text
                        
                        # Also store by slide number for backup lookup
                        slide_key = f"slide_{i+1}"
                        narration_dict[slide_key] = narration_text
                        
                        if not title:
                            print(f"Slide {i+1} (no title): {len(narration_text)} characters of narration")
        
        except Exception as e:
            print(f"Warning: Could not extract narration from slide {i+1}: {e}")
            continue
    
    return narration_dict

# Load the PowerPoint presentation to get slide count and titles
pptx_file = "content_maintenance_process.pptx"
print(f"Loading PowerPoint presentation: {pptx_file}")
presentation = Presentation(pptx_file)
print(f"Loaded presentation with {len(presentation.slides)} slides")

# Extract narration from PowerPoint slide notes
print(f"\nExtracting narration from PowerPoint slide notes...")
narration_notes = extract_narration_from_slides(presentation)

if not narration_notes:
    print("Warning: No narration found in slide notes. Videos will be silent.")
else:
    print(f"Ready to process with {len(narration_notes)} narration entries")

# Try to export slides automatically
print("\nAttempting to export slides as images...")
success = False

# Try PowerShell method first (best for Windows)
print("Trying PowerShell export...")
success = export_slides_as_images_powershell(pptx_file)

# If PowerShell failed, try LibreOffice
if not success:
    print("Trying LibreOffice export...")
    success = export_slides_as_images_libreoffice(pptx_file)

# Look for existing exported slides
slides_dir, slide_files = find_exported_slides()

if not slides_dir:
    print("Could not automatically export slides.")
    export_slides_manual_instruction(pptx_file)
    print("\nAfter manually exporting, run this script again.")
    exit()

print(f"Found {len(slide_files)} exported slide images in {slides_dir}/")
for i, filename in enumerate(slide_files[:5]):  # Show first 5
    print(f"  {i+1}: {filename}")
if len(slide_files) > 5:
    print(f"  ... and {len(slide_files)-5} more")

# Create video from exported slides
video_clips = []
print("\nCreating video from slide images...")

for i, slide_file in enumerate(slide_files):
    slide_path = os.path.join(slides_dir, slide_file)
    print(f"\n--- Processing slide {i+1}/{len(slide_files)}: {slide_file} ---")
    
    try:
        # Load the slide image
        slide_img = Image.open(slide_path)
        print(f"Loaded slide image: {slide_img.size}")
        
        # Resize if needed to standard video dimensions
        # Keep aspect ratio but fit within 1280x720
        slide_img.thumbnail((1280, 720), Image.Resampling.LANCZOS)
        
        # Create a 1280x720 canvas and center the slide
        canvas = Image.new("RGB", (1280, 720), color="white")
        x_offset = (1280 - slide_img.width) // 2
        y_offset = (720 - slide_img.height) // 2
        canvas.paste(slide_img, (x_offset, y_offset))
        
        # Save the processed slide
        processed_path = f"slide_images/processed_slide_{i+1}.png"
        os.makedirs("slide_images", exist_ok=True)
        canvas.save(processed_path)
        
        # Create video clip
        duration = 4  # 4 seconds per slide
        image_clip = ImageClip(processed_path).with_duration(duration)
        video_clips.append(image_clip)
        print(f"Created video clip {i+1} with duration {duration}s")
        
    except Exception as e:
        print(f"Error processing slide {i+1}: {e}")

# Combine clips into final video
if video_clips:
    print(f"\nCombining {len(video_clips)} video clips into final video...")
    final_video = concatenate_videoclips(video_clips, method="compose")
    print("Video clips concatenated successfully")

    output_filename = "code_maintenance_process_SLIDES.mp4"
    print(f"Writing final video file: {output_filename}")
    final_video.write_videofile(output_filename, fps=24)
    print("Final video file written successfully!")

    # Cleanup
    print("Starting cleanup...")
    for clip in video_clips:
        clip.close()
    final_video.close()
    
    print(f"\nVideo created successfully: {output_filename}")
    print(f"Duration: {len(video_clips) * 4} seconds ({len(video_clips)} slides x 4 seconds each)")
else:
    print("No video clips were created!")
